"""
Convert the Beamer presentation PDF into an image-only PPTX for Google Slides.

Google Slides cannot import a PDF, and PDF-to-PPTX converters re-interpret the
LaTeX fonts and break the layout. To keep the Beamer style pixel-perfect, each
PDF page is rasterised with `pdftoppm` and dropped full-bleed onto a blank 4:3
slide, so Slides has nothing to re-flow. The three "final" frames keep their
dashed video placeholder; the actual clips are added by hand in Slides afterwards.

Requires `pdftoppm` (poppler-utils) and the `python-pptx` package.
Usage:
    python docs/presentation/build_pptx.py
    python docs/presentation/build_pptx.py --compile --dpi 300
"""
import argparse
from pathlib import Path
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches

PRESENTATION_DIR = Path(__file__).resolve().parent
DEFAULT_TEX_DIR = PRESENTATION_DIR
DEFAULT_PDF = DEFAULT_TEX_DIR / "build" / "presentation.pdf"
DEFAULT_OUTPUT = PRESENTATION_DIR / "presentation_slides.pptx"


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Build an image-only PPTX from the Beamer PDF.")
    p.add_argument("--pdf", type=Path, default=DEFAULT_PDF,
                   help="Source PDF (default: build/presentation.pdf).")
    p.add_argument("--output", type=Path, default=DEFAULT_OUTPUT,
                   help="Destination PPTX (default: presentation_slides.pptx).")
    p.add_argument("--dpi", type=int, default=300,
                   help="Rasterisation resolution for each page (default: 300).")
    p.add_argument("--compile", action="store_true",
                   help="Recompile the PDF with latexmk before rendering.")
    return p.parse_args()


def compile_pdf(tex_dir: Path) -> None:
    subprocess.run(
        ["latexmk", "-pdf", "-interaction=nonstopmode", "-halt-on-error", "presentation.tex"],
        cwd=tex_dir, check=True,
    )


def render_pages(pdf: Path, dpi: int, out_dir: Path) -> list[Path]:
    # pdftoppm writes <prefix>-NN.png, zero-padded so a plain sort is page order.
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(out_dir / "slide")],
        check=True,
    )
    return sorted(out_dir.glob("slide-*.png"))


def build_pptx(pages: list[Path], output: Path) -> None:
    # Beamer pages are exactly 4:3, so a 10x7.5 in slide matches without distortion.
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(page), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(output))


def main() -> None:
    args = parse_args()

    if args.compile:
        compile_pdf(args.pdf.parent)
    if not args.pdf.exists():
        sys.exit(f"PDF not found: {args.pdf} (run with --compile to build it first).")

    with tempfile.TemporaryDirectory() as tmp:
        pages = render_pages(args.pdf, args.dpi, Path(tmp))
        if not pages:
            sys.exit("pdftoppm produced no pages.")
        build_pptx(pages, args.output)

    size_kb = args.output.stat().st_size // 1024
    print(f"Wrote {len(pages)} slides to {args.output} ({size_kb} KB)")


if __name__ == "__main__":
    main()
